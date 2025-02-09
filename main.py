import streamlit as st
import google.generativeai as genai
import config
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import os



def configure():
    os.environ["GOOGLE_API_KEY"] = config.google_api_key
    genai.configure(api_key=os.environ["GOOGLE_API_KEY"])
    model = genai.GenerativeModel('gemini-pro')
    return model


def generate_ppt_content(topic, num_slides, model):
    prompt = f"""Generate a presentation outline for {topic} with exactly {num_slides} slides.
    Return the response in the following JSON format:
    [
        {{"title": "Title of Slide", "content": ["Point 1", "Point 2", "Point 3"]}}
    ]
    The first slide should be a title slide with a subtitle.
    Each content slide should have 3-4 clear, concise bullet points.
    Make it professional and engaging."""

    try:
        response = model.generate_content(prompt)
        response_text = response.text


        start_idx = response_text.find('[')
        end_idx = response_text.rfind(']') + 1

        if start_idx == -1 or end_idx == 0:
            raise ValueError("Invalid response format")

        json_str = response_text[start_idx:end_idx]
        slides_content = json.loads(json_str)


        if not isinstance(slides_content, list) or not slides_content:
            raise ValueError("Invalid slides content structure")

        return slides_content

    except Exception as e:
        st.error(f"Error generating content: {str(e)}")
        return None


def create_ppt(slides_content):
    prs = Presentation()


    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slides_content[0]['title']
    subtitle.text = slides_content[0]['content'][0] if slides_content[0]['content'] else "Generated with AI"


    for slide_content in slides_content[1:]:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)


        title = slide.shapes.title
        title.text = slide_content['title']


        content = slide.placeholders[1]
        tf = content.text_frame

        for point in slide_content['content']:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0


    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    return output_path


def main():
    st.set_page_config(page_title="AI PPt Generator")

    st.title("AI PPt Generator")
    st.write("Get ppt in minutes")


    try:
        model = configure()
    except Exception as e:
        st.error(f"Error configuring Gemini API: {str(e)}")
        st.stop()


    topic = st.text_input("Enter presentation topic:", placeholder="e.g., Introduction to Artificial Intelligence")
    num_slides = st.slider("Number of slides:", min_value=3, max_value=10, value=5)




    if st.button("Generate Presentation", type="primary"):
        if topic:
            try:
                with st.spinner("in progress.."):

                    slides_content = generate_ppt_content(topic, num_slides, model)

                    if slides_content:

                        output_path = create_ppt(slides_content)


                        with open(output_path, "rb") as file:
                            st.download_button(
                                label="Download ppt",
                                data=file,
                                file_name="generated_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )


                        st.success(" PPt generated successfully:)")
                        st.write("Preview of Generated Content:")
                        for i, slide in enumerate(slides_content):
                            st.write(f"*Slide {i + 1}: {slide['title']}*")
                            for point in slide['content']:
                                st.write(f"- {point}")
                            st.write("---")

            except Exception as e:
                st.error(f"An error occurred: {str(e)}")
        else:
            st.warning("Please enter topic.")


if __name__ == "__main__":
    main()